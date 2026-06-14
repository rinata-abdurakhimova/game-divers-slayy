"""Build the Slay Diver: Rise of 67 pitch deck as a .pptx file.

Run with:
    python docs/pitch/build_deck.py

Regenerates docs/pitch/SlayDiver_Pitch.pptx from the content below. If a
referenced screenshot in docs/pitch/images/ does not exist yet, a labeled
placeholder box is drawn instead so the deck still has the right slide count
and layout.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
IMAGES_DIR = os.path.join(HERE, "images")
OUTPUT = os.path.join(HERE, "SlayDiver_Pitch.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(0x2A, 0x1E, 0x33)
PINK = RGBColor(0xFF, 0x94, 0xD1)
CREAM = RGBColor(0xFF, 0xF2, 0xE8)
PANEL = RGBColor(0x30, 0x24, 0x3D)
ACCENT = RGBColor(0xFF, 0xD1, 0x4D)


def add_background(slide, color=INK):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    # send to back
    sp = rect._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return rect


def add_title(slide, text, top=Inches(0.4), size=40, color=PINK):
    tb = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return tb


def add_bullets(slide, items, left=Inches(0.8), top=Inches(1.6), width=None, height=None,
                 size=22, color=CREAM):
    width = width or (SLIDE_W - Inches(1.6))
    height = height or (SLIDE_H - top - Inches(0.5))
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
    return tb


def add_image_or_placeholder(slide, filename, left, top, width, height, label=None):
    path = os.path.join(IMAGES_DIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width=width, height=height)
        return
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = PINK
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (label or "Screenshot placeholder") + "\n\nAdd file:\ndocs/pitch/images/" + filename
    run.font.size = Pt(16)
    run.font.color.rgb = CREAM
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_box(slide, left, top, width, height, title, body_lines, fill=PANEL, title_color=ACCENT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = PINK
    box.line.width = Pt(1.0)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = title_color
    for line in body_lines:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = line
        run2.font.size = Pt(12)
        run2.font.color.rgb = CREAM
    return box


def add_connector(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    conn.line.color.rgb = PINK
    conn.line.width = Pt(1.5)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "Slay Diver: Rise of 67", top=Inches(2.4), size=54)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(3.6), SLIDE_W - Inches(1.2), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "A boss-chase platformer where the only weapon is arithmetic."
    run.font.size = Pt(26)
    run.font.color.rgb = CREAM
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Play it now: https://alina-anila.itch.io/slay-diver-rise-of-67"
    run2.font.size = Pt(20)
    run2.font.color.rgb = ACCENT
    p2.space_before = Pt(20)

    # 2. Hook
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "The Hook")
    quote = s.shapes.add_textbox(Inches(0.8), Inches(1.4), SLIDE_W - Inches(1.6), Inches(1.0))
    qp = quote.text_frame.paragraphs[0]
    qrun = qp.add_run()
    qrun.text = "Make the score exactly 67 while Boss 67 tries to mathematically ruin you."
    qrun.font.italic = True
    qrun.font.size = Pt(28)
    qrun.font.color.rgb = ACCENT
    add_bullets(s, [
        "Score starts at 1.00",
        "Land on exactly 67.00 → you win",
        "Land on 0.00 → the run fails",
        "Everything on screen is just another +, -, or x on your score",
    ], top=Inches(2.6))

    # 3. Land phase
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "Gameflow: Land Phase", size=34)
    add_image_or_placeholder(s, "01-land-phase.png", Inches(0.6), Inches(1.4),
                              Inches(7.6), Inches(5.4), "Land phase gameplay")
    add_bullets(s, [
        "Walk, jump, and dodge across the authored sand route",
        "Green pickups add to your score (+1...+7)",
        "White boss numbers multiply your score (x0, x0.5, x0.8) and are blocked by terrain",
        "HUD tracks score, target 67, distance in blocks, and recent operations",
    ], left=Inches(8.4), top=Inches(1.6), width=Inches(4.5), size=18)

    # 4. Water reversed
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "The Blue Flood: Reversed Controls", size=32)
    add_image_or_placeholder(s, "02-water-reversed.png", Inches(0.6), Inches(1.4),
                              Inches(7.6), Inches(5.4), "Water event - reversed controls")
    add_bullets(s, [
        "At 28 blocks (or when score is divisible by 6/7), the Blue Flood hits",
        "Water rules swap operations: boss subtracts, floor pickups multiply",
        "50/50 complication roll: here, left/right controls are reversed for 20s",
    ], left=Inches(8.4), top=Inches(1.6), width=Inches(4.5), size=18)

    # 5. Water gravity
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "The Blue Flood: Gravity Inverted", size=32)
    add_image_or_placeholder(s, "03-water-gravity.png", Inches(0.6), Inches(1.4),
                              Inches(7.6), Inches(5.4), "Water event - gravity inverted")
    add_bullets(s, [
        "Other 50/50 outcome: the whole world flips 180°",
        "Camera and player render upside down, sand rises to the top",
        "HUD panels swap top/bottom so score and rules stay readable",
        "Physics and operand logic are untouched — purely visual disorientation",
    ], left=Inches(8.4), top=Inches(1.6), width=Inches(4.5), size=18)

    # 6. Game over
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "Gameflow: Failure State", size=34)
    add_image_or_placeholder(s, "04-game-over.png", Inches(0.6), Inches(1.4),
                              Inches(7.6), Inches(5.4), "Game over screen")
    add_bullets(s, [
        "Score hits 0 → \"THE NUMBERS WON\"",
        "Recent-operations panel shows exactly which hits caused the fail",
        "Instant restart with R — no friction between attempts",
    ], left=Inches(8.4), top=Inches(1.6), width=Inches(4.5), size=18)

    # 7. Architecture overview
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "Architecture Overview", size=34)

    col_w = Inches(2.9)
    col_h = Inches(1.4)
    gap = Inches(0.3)
    top0 = Inches(1.5)

    # Autoloads column
    add_box(s, Inches(0.6), top0, col_w, col_h, "GameState",
            ["score, run phase, water state,", "win / fail"])
    add_box(s, Inches(0.6), top0 + col_h + gap, col_w, col_h, "GameEvents",
            ["global signal hub"])
    add_box(s, Inches(0.6), top0 + 2 * (col_h + gap), col_w, col_h, "AudioBus / SceneLoader",
            ["sound playback, intro,", "restart flow"])

    # Shared rules column
    add_box(s, Inches(0.6) + col_w + gap, top0, col_w, col_h, "GameRules",
            ["constants & operations"])
    add_box(s, Inches(0.6) + col_w + gap, top0 + col_h + gap, col_w, col_h, "ScoreService",
            ["fixed-point score math"])
    add_box(s, Inches(0.6) + col_w + gap, top0 + 2 * (col_h + gap), col_w, col_h, "WaterRuleService",
            ["water variant &", "complication selection"])

    # World column
    add_box(s, Inches(0.6) + 2 * (col_w + gap), top0, col_w, col_h, "Boss67LevelController",
            ["chunks, spawns, distance,", "water triggers"])
    add_box(s, Inches(0.6) + 2 * (col_w + gap), top0 + col_h + gap, col_w, col_h, "Player",
            ["platformer movement,", "phase visuals"])
    add_box(s, Inches(0.6) + 2 * (col_w + gap), top0 + 2 * (col_h + gap), col_w, col_h, "Boss67",
            ["presentation &", "projectile spawn anchors"])

    # UI column
    add_box(s, Inches(0.6) + 3 * (col_w + gap), top0, col_w, col_h, "HUD",
            ["score, rules, water timer,", "power-ups"])
    add_box(s, Inches(0.6) + 3 * (col_w + gap), top0 + col_h + gap, col_w, col_h, "ResultScreen",
            ["victory / failure, restart"])
    add_box(s, Inches(0.6) + 3 * (col_w + gap), top0 + 2 * (col_h + gap), col_w, col_h,
            "WaterRuleOverlay", ["water rule announcement card"])

    notes = s.shapes.add_textbox(Inches(0.6), top0 + 3 * (col_h + gap) + Inches(0.1),
                                  SLIDE_W - Inches(1.2), Inches(1.2))
    tf = notes.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("Contract-first: every cross-owner signal lives in GameEvents and is documented in "
                 "INTEGRATION_CONTRACT.md. Three ownership zones (Polina / Alina / Rinata) work in "
                 "parallel without touching each other's scenes directly.")
    run.font.size = Pt(16)
    run.font.color.rgb = ACCENT
    run.font.italic = True

    # 8. Our process
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "Our Process", size=34)
    add_image_or_placeholder(s, "05-blackboard.png", Inches(0.6), Inches(1.4),
                              Inches(7.6), Inches(5.4), "Whiteboard planning session")
    add_bullets(s, [
        "Started with a chalkboard jam session: cutscene plan, Boss 67 mechanics, "
        "water variants, operand tables, and the authored level layout",
        "Locked the numeric rules (x0/x0.5/x0.8, water A/B/C variants, +1..+7 pickups) "
        "before writing code",
        "Split into three ownership zones from day one, integrated continuously via a "
        "shared contract doc",
    ], left=Inches(8.4), top=Inches(1.6), width=Inches(4.5), size=18)

    # 9. Closing
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_title(s, "Try It", top=Inches(2.6), size=44)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(3.8), SLIDE_W - Inches(1.2), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "https://alina-anila.itch.io/slay-diver-rise-of-67"
    run.font.size = Pt(28)
    run.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Avoid zero. Use your operations. Hit exactly 67."
    run2.font.size = Pt(22)
    run2.font.color.rgb = CREAM
    p2.space_before = Pt(16)

    prs.save(OUTPUT)
    print("Saved", OUTPUT)


if __name__ == "__main__":
    build()
